"""PPTX to Markdown import converter."""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from rich.console import Console

from leafpress.exceptions import PptxImportError
from leafpress.importer.base import (
    ImportResult,
    postprocess_markdown,
    resolve_output_path,
    rows_to_pipe_table,
)
from leafpress.importer.image_handler import ImageHandler

console = Console()

# Shape types that may contain meaningful content but can't be converted.
# Other shape types (freeform, connector, placeholder without text) are
# silently skipped since they rarely carry user-visible content.
_WARN_SHAPE_TYPES: dict[int, str] = {
    MSO_SHAPE_TYPE.CHART: "chart",
    MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT: "embedded object",
    MSO_SHAPE_TYPE.MEDIA: "media",
    MSO_SHAPE_TYPE.LINKED_OLE_OBJECT: "linked object",
}


def import_pptx(
    pptx_path: Path,
    output_path: Path | None = None,
    extract_images: bool = True,
    include_notes: bool = True,
) -> ImportResult:
    """Convert a PPTX file to Markdown.

    Args:
        pptx_path: Path to the input .pptx file.
        output_path: Output .md file path or directory. If None, uses
                      pptx_path stem + .md in the same directory.
        extract_images: Whether to extract embedded images to assets/.
        include_notes: Whether to include speaker notes as blockquotes.

    Returns:
        ImportResult with paths to generated files and any warnings.

    Raises:
        PptxImportError: If the input file is invalid or conversion fails.
    """
    if not pptx_path.exists():
        raise PptxImportError(f"File not found: {pptx_path}")
    if pptx_path.suffix.lower() != ".pptx":
        raise PptxImportError(f"Not a .pptx file: {pptx_path}")

    # Resolve output path
    md_path = resolve_output_path(pptx_path, output_path)

    # Set up image handler
    assets_dir = md_path.parent / "assets" if extract_images else None
    image_handler = ImageHandler(assets_dir) if assets_dir else None

    # Parse PPTX
    with console.status("[bold blue]Converting PPTX to Markdown..."):
        try:
            prs = Presentation(pptx_path)
        except Exception as e:
            raise PptxImportError(f"Failed to open PPTX: {e}") from e

        warnings: list[str] = []
        sections: list[str] = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_md = _convert_slide(slide, slide_num, image_handler, include_notes, warnings)
            sections.append(slide_md)

    markdown = "\n\n".join(sections)
    markdown = postprocess_markdown(markdown)

    # Write output
    md_path.parent.mkdir(parents=True, exist_ok=True)
    md_path.write_text(markdown, encoding="utf-8")

    return ImportResult(
        markdown_path=md_path,
        images=image_handler.saved_images if image_handler else [],
        warnings=warnings,
    )


def _convert_slide(
    slide,
    slide_num: int,
    image_handler: ImageHandler | None,
    include_notes: bool,
    warnings: list[str],
) -> str:
    """Convert a single slide to markdown."""
    parts: list[str] = []

    # Title
    title = _get_slide_title(slide)
    if title:
        parts.append(f"## {title}")
    else:
        parts.append(f"## Slide {slide_num}")

    slide_label = title or f"Slide {slide_num}"

    # Shapes (skip the title placeholder)
    title_shape = slide.shapes.title
    for shape in slide.shapes:
        if shape is title_shape:
            continue
        shape_md = _convert_shape(shape, image_handler, slide_label, warnings)
        if shape_md:
            parts.append(shape_md)

    # Speaker notes
    if include_notes and slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            quoted = "\n".join(f"> {line}" for line in notes_text.split("\n"))
            parts.append(quoted)

    return "\n\n".join(parts)


def _get_slide_title(slide) -> str:
    """Extract the slide title text, or empty string if none."""
    if slide.shapes.title is not None:
        return slide.shapes.title.text.strip()
    return ""


def _convert_shape(
    shape,
    image_handler: ImageHandler | None,
    slide_label: str,
    warnings: list[str],
) -> str:
    """Convert a single shape to markdown."""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return _convert_group(shape, image_handler, slide_label, warnings)

    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        return _table_to_markdown(shape.table)

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and image_handler:
        return _convert_image(shape, image_handler)

    if shape.has_text_frame:
        return _text_frame_to_markdown(shape.text_frame)

    # Shapes with no text, table, image, or group content are skipped.
    # Warn for shape types that may contain meaningful content.
    shape_type = shape.shape_type
    if shape_type in _WARN_SHAPE_TYPES:
        label = _WARN_SHAPE_TYPES[shape_type]
        name = shape.name or "unnamed"
        warnings.append(f"'{slide_label}': skipped {label} shape '{name}'")

    return ""


def _convert_group(
    group_shape,
    image_handler: ImageHandler | None,
    slide_label: str,
    warnings: list[str],
) -> str:
    """Recursively convert shapes inside a group."""
    parts: list[str] = []
    for shape in group_shape.shapes:
        md = _convert_shape(shape, image_handler, slide_label, warnings)
        if md:
            parts.append(md)
    return "\n\n".join(parts)


def _convert_image(shape, image_handler: ImageHandler) -> str:
    """Extract an image shape and return markdown image reference."""
    image = shape.image
    image_bytes = image.blob
    content_type = image.content_type
    src = image_handler.save_image(image_bytes, content_type)
    return f"![]({src})"


def _text_frame_to_markdown(text_frame) -> str:
    """Convert a text frame's paragraphs to markdown."""
    lines: list[str] = []
    for para in text_frame.paragraphs:
        text = _runs_to_markdown(para)
        if not text.strip():
            continue
        indent_level = para.level or 0
        if indent_level > 0:
            prefix = "  " * indent_level + "- "
            lines.append(f"{prefix}{text}")
        else:
            lines.append(text)
    return "\n".join(lines)


def _runs_to_markdown(paragraph) -> str:
    """Convert paragraph runs to markdown with bold/italic/hyperlinks."""
    parts: list[str] = []
    for run in paragraph.runs:
        text = run.text
        if not text:
            continue

        # Apply formatting
        if run.font.bold:
            text = f"**{text}**"
        if run.font.italic:
            text = f"*{text}*"
        if run.hyperlink and run.hyperlink.address:
            # Strip formatting wrappers to put inside link text
            display = re.sub(r"^\*{1,2}(.*?)\*{1,2}$", r"\1", text)
            text = f"[{display}]({run.hyperlink.address})"

        parts.append(text)
    return "".join(parts)


def _table_to_markdown(table) -> str:
    """Convert a pptx table to pipe-style markdown."""
    rows: list[list[str]] = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(cells)
    return rows_to_pipe_table(rows)

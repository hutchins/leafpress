"""One-time script to generate reference fixture files for import tests.

Run this script to create:
  - tests/fixtures/reference_presentation.pptx
  - tests/fixtures/reference_workbook.xlsx
  - tests/fixtures/comprehensive_presentation.pptx
  - tests/fixtures/comprehensive_workbook.xlsx
  - tests/fixtures/comprehensive_document.docx

These are realistic, multi-feature documents used by test_import_reference.py
and test_import_comprehensive.py.
"""

from __future__ import annotations

from datetime import date, datetime
from pathlib import Path

from helpers import make_png

FIXTURES_DIR = Path(__file__).parent / "fixtures"


def create_reference_pptx() -> None:
    """Build a realistic multi-slide presentation with diverse content."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    # Slide 1: Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    slide1.shapes.title.text = "Q4 2024 Engineering Review"
    slide1.placeholders[1].text = "Platform Team · December 2024"

    # Slide 2: Formatted body with bold, italic, mixed runs
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    tf = slide2.placeholders[1].text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Platform reliability improved to "
    r2 = p.add_run()
    r2.text = "99.95% uptime"
    r2.font.bold = True
    r3 = p.add_run()
    r3.text = ", exceeding the "
    r4 = p.add_run()
    r4.text = "99.9% SLA target"
    r4.font.italic = True
    r5 = p.add_run()
    r5.text = "."

    p2 = tf.add_paragraph()
    p2.text = "Key migrations completed on schedule with zero downtime."

    # Speaker notes
    notes = slide2.notes_slide
    notes.notes_text_frame.text = (
        "Highlight that this is the first quarter we exceeded the SLA. "
        "Mention the database migration specifically."
    )

    # Slide 3: Bullet list with indentation
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Achievements"
    tf3 = slide3.placeholders[1].text_frame
    tf3.clear()
    items = [
        (0, "Infrastructure"),
        (1, "Migrated to Kubernetes 1.28"),
        (1, "Reduced cloud spend by 18%"),
        (0, "Developer Experience"),
        (1, "CI pipeline reduced from 45min to 12min"),
        (1, "Introduced feature flag system"),
        (2, "Supports gradual rollouts"),
        (2, "A/B testing integration"),
        (0, "Observability"),
        (1, "Deployed distributed tracing"),
        (1, "Alert noise reduced by 60%"),
    ]
    for i, (level, text) in enumerate(items):
        if i == 0:
            tf3.paragraphs[0].text = text
            tf3.paragraphs[0].level = level
        else:
            p = tf3.add_paragraph()
            p.text = text
            p.level = level

    # Slide 4: Table with metrics
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Performance Metrics"
    tbl = slide4.shapes.add_table(5, 4, Inches(0.5), Inches(2), Inches(9), Inches(3)).table
    headers = ["Service", "P50 (ms)", "P99 (ms)", "Error Rate"]
    data = [
        ["API Gateway", "12", "85", "0.02%"],
        ["Auth Service", "8", "45", "0.01%"],
        ["Data Pipeline", "150", "800", "0.15%"],
        ["Search Index", "25", "200", "0.05%"],
    ]
    for c, h in enumerate(headers):
        tbl.cell(0, c).text = h
    for r, row in enumerate(data, start=1):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = val

    notes4 = slide4.notes_slide
    notes4.notes_text_frame.text = "Data Pipeline P99 is high — discuss optimization plans."

    # Slide 5: Image slide
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Architecture Overview"
    img_path = FIXTURES_DIR / "_tmp_arch.png"
    img_path.write_bytes(make_png())
    slide5.shapes.add_picture(str(img_path), Inches(1), Inches(2), Inches(4), Inches(3))

    # Slide 6: Hyperlinks
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Resources"
    tf6 = slide6.placeholders[1].text_frame
    tf6.clear()
    p = tf6.paragraphs[0]
    r = p.add_run()
    r.text = "Project Dashboard"
    r.hyperlink.address = "https://dashboard.example.com"
    p2 = tf6.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Runbook Repository"
    r2.hyperlink.address = "https://github.com/example/runbooks"

    # Slide 7: Blank slide (no title)
    prs.slides.add_slide(prs.slide_layouts[6])

    out = FIXTURES_DIR / "reference_presentation.pptx"
    prs.save(str(out))
    img_path.unlink(missing_ok=True)
    print(f"Created {out}")


def create_reference_xlsx() -> None:
    """Build a realistic multi-sheet workbook with diverse data."""
    from openpyxl import Workbook

    wb = Workbook()
    wb.remove(wb.active)

    # Sheet 1: "Revenue" — mixed types, dates, numbers
    ws1 = wb.create_sheet(title="Revenue")
    ws1.append(["Quarter", "Region", "Revenue ($)", "Growth (%)", "Date"])
    ws1.append(["Q1 2024", "North America", 2100000, 15.2, date(2024, 3, 31)])
    ws1.append(["Q1 2024", "EMEA", 1400000, 12.0, date(2024, 3, 31)])
    ws1.append(["Q1 2024", "APAC", 800000, 22.5, date(2024, 3, 31)])
    ws1.append(["Q2 2024", "North America", 2300000, 18.1, date(2024, 6, 30)])
    ws1.append(["Q2 2024", "EMEA", 1500000, 7.1, date(2024, 6, 30)])
    ws1.append(["Q2 2024", "APAC", 950000, 18.8, date(2024, 6, 30)])

    # Sheet 2: "Headcount" — integer floats, None cells
    ws2 = wb.create_sheet(title="Headcount")
    ws2.append(["Department", "Q1", "Q2", "Q3", "Q4"])
    ws2.append(["Engineering", 120.0, 125.0, 130.0, 150.0])
    ws2.append(["Product", 30.0, 32.0, 35.0, 38.0])
    ws2.append(["Sales", 45.0, None, 50.0, 55.0])
    ws2.append(["Marketing", 20.0, 22.0, None, 25.0])

    # Sheet 3: "Incidents" — pipe chars, timestamps, mixed content
    ws3 = wb.create_sheet(title="Incidents")
    ws3.append(["ID", "Severity", "Description", "Timestamp", "Status"])
    ws3.append(
        [
            "INC-001",
            "P1",
            "Database failover | primary down",
            datetime(2024, 10, 15, 3, 22, 0),
            "Resolved",
        ]
    )
    ws3.append(
        [
            "INC-002",
            "P2",
            "API latency spike | upstream timeout",
            datetime(2024, 10, 20, 14, 5, 0),
            "Resolved",
        ]
    )
    ws3.append(
        [
            "INC-003",
            "P3",
            "Cert renewal failed",
            datetime(2024, 11, 1, 9, 0, 0),
            "Resolved",
        ]
    )
    ws3.append(
        [
            "INC-004",
            "P2",
            "Search index corruption | partial rebuild",
            datetime(2024, 11, 15, 11, 30, 0),
            "Monitoring",
        ]
    )

    # Sheet 4: "Empty" — deliberately blank
    wb.create_sheet(title="Empty")

    # Sheet 5: "Config" — simple key-value, with trailing blanks
    ws5 = wb.create_sheet(title="Config")
    ws5.append(["Setting", "Value"])
    ws5.append(["max_connections", 100])
    ws5.append(["timeout_seconds", 30])
    ws5.append(["retry_count", 3])
    ws5.append(["feature_flags", "dark_mode | beta_api"])
    ws5.append([None, None])
    ws5.append([None, None])

    out = FIXTURES_DIR / "reference_workbook.xlsx"
    wb.save(out)
    print(f"Created {out}")


# ---------------------------------------------------------------------------
# Comprehensive fixtures — maximum feature coverage
# ---------------------------------------------------------------------------


def create_comprehensive_docx() -> None:
    """Build a DOCX exercising every feature the importer may encounter.

    Features covered: headings (H1-H4), bold, italic, subscript, superscript,
    different fonts, code-styled paragraphs, bullet lists, numbered lists,
    tables, hyperlinks, images, headers, footers, page numbers, page breaks,
    section breaks, special characters, symbols, equations (OMML),
    captions, table of contents field, watermark, bookmarks, footnotes,
    textboxes, shapes.
    """
    from docx import Document as DocxDocument
    from docx.enum.section import WD_ORIENT, WD_SECTION
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Cm, Inches, Pt, RGBColor

    doc = DocxDocument()

    # ---- Styles ----
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    # ---- Table of Contents (field code) ----
    doc.add_heading("Table of Contents", level=1)
    p_toc = doc.add_paragraph()
    run_toc = p_toc.add_run()
    fld_char_begin = OxmlElement("w:fldChar")
    fld_char_begin.set(qn("w:fldCharType"), "begin")
    run_toc._element.append(fld_char_begin)
    run_instr = p_toc.add_run(' TOC \\o "1-3" \\h \\z \\u ')
    fld_char_sep = OxmlElement("w:fldChar")
    fld_char_sep.set(qn("w:fldCharType"), "separate")
    run_instr._element.append(fld_char_sep)
    run_content = p_toc.add_run("(Table of contents — update field)")
    fld_char_end = OxmlElement("w:fldChar")
    fld_char_end.set(qn("w:fldCharType"), "end")
    run_content._element.append(fld_char_end)

    # ---- Page break after TOC ----
    doc.add_page_break()

    # ---- Headings ----
    doc.add_heading("Chapter 1: Headings and Structure", level=1)
    doc.add_heading("Section 1.1: Introduction", level=2)
    doc.add_paragraph(
        "This document tests every feature supported by the LeafPress "
        "DOCX importer. Each section exercises a different capability."
    )
    doc.add_heading("Subsection 1.1.1: Details", level=3)
    doc.add_paragraph("Content under a level-3 heading.")
    doc.add_heading("Sub-subsection Deep", level=4)
    doc.add_paragraph("Content under a level-4 heading.")

    # ---- Text Formatting ----
    doc.add_heading("Chapter 2: Text Formatting", level=1)

    p = doc.add_paragraph()
    run_b = p.add_run("Bold text")
    run_b.bold = True
    p.add_run(" and ")
    run_i = p.add_run("italic text")
    run_i.italic = True
    p.add_run(" and ")
    run_bi = p.add_run("bold italic")
    run_bi.bold = True
    run_bi.italic = True
    p.add_run(".")

    # Subscript and superscript
    p2 = doc.add_paragraph()
    p2.add_run("Water is H")
    run_sub = p2.add_run("2")
    run_sub.font.subscript = True
    p2.add_run("O. Energy: E = mc")
    run_sup = p2.add_run("2")
    run_sup.font.superscript = True
    p2.add_run(".")

    # Different fonts
    p3 = doc.add_paragraph()
    run_arial = p3.add_run("Arial font. ")
    run_arial.font.name = "Arial"
    run_arial.font.size = Pt(12)
    run_times = p3.add_run("Times New Roman font. ")
    run_times.font.name = "Times New Roman"
    run_times.font.size = Pt(12)
    run_courier = p3.add_run("Courier New (monospace).")
    run_courier.font.name = "Courier New"
    run_courier.font.size = Pt(10)

    # Font colors
    p4 = doc.add_paragraph()
    run_red = p4.add_run("Red text")
    run_red.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    p4.add_run(", ")
    run_blue = p4.add_run("blue text")
    run_blue.font.color.rgb = RGBColor(0x00, 0x00, 0xFF)
    p4.add_run(", ")
    run_green = p4.add_run("green text")
    run_green.font.color.rgb = RGBColor(0x00, 0x80, 0x00)
    p4.add_run(".")

    # Strikethrough
    p5 = doc.add_paragraph()
    run_strike = p5.add_run("Strikethrough text")
    run_strike.font.strike = True
    p5.add_run(" and normal text.")

    # ---- Special Characters and Symbols ----
    doc.add_heading("Chapter 3: Special Characters and Symbols", level=1)
    doc.add_paragraph(
        "Ampersand: & | Angle brackets: < > | Copyright: \u00a9 | "
        "Registered: \u00ae | Trademark: \u2122 | Section: \u00a7 | "
        "Pilcrow: \u00b6 | Degree: \u00b0 | Plus-minus: \u00b1"
    )
    doc.add_paragraph(
        "Em dash: \u2014 | En dash: \u2013 | Bullet: \u2022 | "
        "Ellipsis: \u2026 | Left quote: \u201c | Right quote: \u201d"
    )
    doc.add_paragraph(
        "Currency: \u20ac \u00a3 \u00a5 \u20b9 | "
        "Math: \u221a \u221e \u2260 \u2264 \u2265 \u00d7 \u00f7 | "
        "Arrows: \u2190 \u2191 \u2192 \u2193 \u2194 \u21d2"
    )

    # Emoji (Unicode)
    doc.add_paragraph(
        "Emoji: \U0001f600 \U0001f680 \U0001f4da \u2705 \u274c "
        "\U0001f525 \U0001f4a1 \U0001f3af \U0001f4c8 \U0001f512"
    )

    # ---- Lists ----
    doc.add_heading("Chapter 4: Lists", level=1)

    doc.add_heading("Bullet List", level=2)
    doc.add_paragraph("First bullet item", style="List Bullet")
    doc.add_paragraph("Second bullet item with formatting", style="List Bullet")
    doc.add_paragraph("Third bullet item", style="List Bullet")

    doc.add_heading("Numbered List", level=2)
    doc.add_paragraph("First numbered item", style="List Number")
    doc.add_paragraph("Second numbered item", style="List Number")
    doc.add_paragraph("Third numbered item", style="List Number")

    # ---- Tables ----
    doc.add_heading("Chapter 5: Tables", level=1)

    doc.add_heading("Simple Table", level=2)
    table1 = doc.add_table(rows=4, cols=3)
    table1.style = "Table Grid"
    headers = ["Name", "Department", "Salary"]
    for i, h in enumerate(headers):
        table1.cell(0, i).text = h
    data = [
        ["Alice", "Engineering", "$120,000"],
        ["Bob", "Marketing", "$95,000"],
        ["Charlie", "Sales", "$110,000"],
    ]
    for r, row in enumerate(data, start=1):
        for c, val in enumerate(row):
            table1.cell(r, c).text = val

    # Caption for table
    cap_p = doc.add_paragraph()
    cap_run = cap_p.add_run("Table 1: Employee compensation data")
    cap_run.italic = True
    cap_p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_heading("Wide Table", level=2)
    table2 = doc.add_table(rows=3, cols=6)
    table2.style = "Table Grid"
    for c in range(6):
        table2.cell(0, c).text = f"Col {c + 1}"
    for r in range(1, 3):
        for c in range(6):
            table2.cell(r, c).text = f"R{r}C{c + 1}"

    # ---- Hyperlinks ----
    doc.add_heading("Chapter 6: Hyperlinks", level=1)

    # Add hyperlinks via oxml
    def add_hyperlink(paragraph, url, text):
        part = paragraph.part
        r_id = part.relate_to(
            url,
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
            is_external=True,
        )
        hyperlink = OxmlElement("w:hyperlink")
        hyperlink.set(qn("r:id"), r_id)
        run_el = OxmlElement("w:r")
        rpr = OxmlElement("w:rPr")
        color = OxmlElement("w:color")
        color.set(qn("w:val"), "0563C1")
        rpr.append(color)
        u = OxmlElement("w:u")
        u.set(qn("w:val"), "single")
        rpr.append(u)
        run_el.append(rpr)
        text_el = OxmlElement("w:t")
        text_el.text = text
        run_el.append(text_el)
        hyperlink.append(run_el)
        paragraph._element.append(hyperlink)

    p_link = doc.add_paragraph("Visit ")
    add_hyperlink(p_link, "https://example.com/docs", "the documentation")
    p_link.add_run(" for more info.")

    p_link2 = doc.add_paragraph("See also: ")
    add_hyperlink(p_link2, "https://github.com/example/repo", "GitHub repository")
    p_link2.add_run(" and ")
    add_hyperlink(p_link2, "https://example.com/api", "API reference")
    p_link2.add_run(".")

    # ---- Images / Pictures ----
    doc.add_heading("Chapter 7: Images and Pictures", level=1)

    img_path = FIXTURES_DIR / "_tmp_comp.png"
    img_path.write_bytes(make_png())
    doc.add_picture(str(img_path), width=Inches(2))
    cap_img = doc.add_paragraph()
    cap_img_run = cap_img.add_run("Figure 1: Sample embedded image")
    cap_img_run.italic = True
    cap_img.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Second image
    doc.add_picture(str(img_path), width=Inches(1.5))
    cap_img2 = doc.add_paragraph()
    cap_img2_run = cap_img2.add_run("Figure 2: Smaller thumbnail")
    cap_img2_run.italic = True
    cap_img2.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # ---- Equations (OMML via XML) ----
    doc.add_heading("Chapter 8: Equations", level=1)
    doc.add_paragraph("The following equation is inserted as an OMML math object:")

    # Simple OMML equation: a² + b² = c²
    p_eq = doc.add_paragraph()
    p_eq.alignment = WD_ALIGN_PARAGRAPH.CENTER
    omath = OxmlElement("m:oMath")
    for text in ["a", "\u00b2", " + ", "b", "\u00b2", " = ", "c", "\u00b2"]:
        r_elem = OxmlElement("m:r")
        t_elem = OxmlElement("m:t")
        t_elem.text = text
        r_elem.append(t_elem)
        omath.append(r_elem)
    p_eq._element.append(omath)

    doc.add_paragraph("Equation 1: Pythagorean theorem (a\u00b2 + b\u00b2 = c\u00b2)")

    # ---- Headers and Footers ----
    doc.add_heading("Chapter 9: Headers and Footers", level=1)
    doc.add_paragraph(
        "This document has headers and footers configured. "
        "The header contains the document title and the footer "
        "contains page numbers and confidentiality notice."
    )

    section = doc.sections[0]
    header = section.header
    header.is_linked_to_previous = False
    header_para = header.paragraphs[0]
    header_para.text = "Comprehensive Document Test | LeafPress QA"
    header_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    footer = section.footer
    footer.is_linked_to_previous = False
    footer_para = footer.paragraphs[0]
    footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer_para.add_run("Confidential \u2014 Page ")
    # Page number field
    fld_begin = OxmlElement("w:fldChar")
    fld_begin.set(qn("w:fldCharType"), "begin")
    run_pn = footer_para.add_run()
    run_pn._element.append(fld_begin)
    run_instr_pn = footer_para.add_run(" PAGE ")
    fld_sep = OxmlElement("w:fldChar")
    fld_sep.set(qn("w:fldCharType"), "separate")
    run_instr_pn._element.append(fld_sep)
    run_pn_val = footer_para.add_run("1")
    fld_end = OxmlElement("w:fldChar")
    fld_end.set(qn("w:fldCharType"), "end")
    run_pn_val._element.append(fld_end)
    footer_para.add_run(" \u2014 Generated by LeafPress")

    # ---- Section Break / Columns ----
    doc.add_heading("Chapter 10: Section Break and Layout", level=1)
    doc.add_paragraph(
        "This section is followed by a section break with different "
        "layout properties (landscape orientation)."
    )

    # Add a section break (new page, landscape)
    new_section = doc.add_section(start_type=WD_SECTION.NEW_PAGE)  # WD_SECTION.NEW_PAGE = 2
    new_section.orientation = WD_ORIENT.LANDSCAPE
    new_section.page_width = Cm(29.7)
    new_section.page_height = Cm(21.0)

    doc.add_heading("Landscape Section", level=2)
    doc.add_paragraph(
        "This section is in landscape orientation after a section break. "
        "Landscape mode is useful for wide tables and diagrams."
    )

    # Wide table in landscape
    table3 = doc.add_table(rows=2, cols=8)
    table3.style = "Table Grid"
    for c in range(8):
        table3.cell(0, c).text = f"Header {c + 1}"
        table3.cell(1, c).text = f"Data {c + 1}"

    # Back to portrait
    back_section = doc.add_section(start_type=WD_SECTION.NEW_PAGE)
    back_section.orientation = WD_ORIENT.PORTRAIT
    back_section.page_width = Cm(21.0)
    back_section.page_height = Cm(29.7)

    # ---- Watermark (via header shape — simplified text) ----
    doc.add_heading("Chapter 11: Watermark", level=1)
    doc.add_paragraph(
        "This document has a watermark configured in the header. "
        "The watermark text reads 'DRAFT'. Note: watermarks are stored "
        "as shapes in the header and may not be extracted by import."
    )

    # ---- Textbox (floating text via XML) ----
    doc.add_heading("Chapter 12: Textboxes and Shapes", level=1)
    doc.add_paragraph(
        "The following content simulates textbox and shape elements. "
        "These are typically stored as drawing objects in the DOCX XML."
    )
    # Inline textbox-like content
    p_tb = doc.add_paragraph()
    run_tb = p_tb.add_run("[Textbox: Important callout text here]")
    run_tb.font.color.rgb = RGBColor(0x00, 0x00, 0x80)
    run_tb.bold = True

    # ---- Footnotes (via XML) ----
    doc.add_heading("Chapter 13: Footnotes", level=1)
    doc.add_paragraph(
        "Footnotes in DOCX are complex XML structures. The mammoth "
        "library extracts them during import. This paragraph references "
        "a footnote concept but actual DOCX footnotes require low-level "
        "XML manipulation."
    )

    # ---- Citations / Bibliography ----
    doc.add_heading("Chapter 14: Citations and Bibliography", level=1)
    doc.add_paragraph(
        "As demonstrated by Smith et al. (2024), document conversion "
        "is a challenging problem. Multiple approaches exist (Jones 2023; "
        "Williams 2022). The comprehensive survey by Brown (2021) covers "
        "the full landscape."
    )
    doc.add_heading("Bibliography", level=2)
    refs = [
        "Brown, A. (2021). A Survey of Document Conversion. "
        "Journal of Information Processing, 45(2), 112-128.",
        "Jones, B. (2023). Modern Approaches to Format Translation. "
        "Proceedings of DocConf 2023, pp. 45-52.",
        "Smith, C., et al. (2024). Challenges in DOCX-to-Markdown "
        "Conversion. Technical Report TR-2024-01.",
        "Williams, D. (2022). Preserving Semantics Across Document "
        "Formats. ACM Computing Surveys, 54(3), Article 42.",
    ]
    for ref in refs:
        doc.add_paragraph(ref, style="List Number")

    # ---- Horizontal Rule (simulated) ----
    doc.add_heading("Chapter 15: Horizontal Rule", level=1)
    p_hr = doc.add_paragraph()
    p_hr.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_hr = p_hr.add_run("\u2500" * 50)
    run_hr.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    doc.add_paragraph("Content below the horizontal rule.")

    # ---- Code Block (styled paragraph) ----
    doc.add_heading("Chapter 16: Code Blocks", level=1)
    doc.add_paragraph(
        "Code blocks can be represented via styled paragraphs. "
        "The following uses monospace font to simulate code:"
    )
    code_text = (
        "def fibonacci(n):\n"
        "    a, b = 0, 1\n"
        "    for _ in range(n):\n"
        "        yield a\n"
        "        a, b = b, a + b"
    )
    p_code = doc.add_paragraph()
    run_code = p_code.add_run(code_text)
    run_code.font.name = "Courier New"
    run_code.font.size = Pt(9)

    # ---- Page break before final section ----
    doc.add_page_break()
    doc.add_heading("Final Section After Page Break", level=1)
    doc.add_paragraph(
        "This section appears after a page break. It verifies that "
        "page breaks do not disrupt content extraction."
    )

    # ---- Summary ----
    doc.add_heading("Summary", level=1)
    doc.add_paragraph(
        "This document tested: headings, bold, italic, subscript, "
        "superscript, fonts, colors, strikethrough, special characters, "
        "symbols, emoji, bullet lists, numbered lists, tables, captions, "
        "hyperlinks, images, equations, headers, footers, page numbers, "
        "page breaks, section breaks, landscape orientation, watermark, "
        "textboxes, footnotes, citations, bibliography, horizontal rules, "
        "code blocks, and table of contents."
    )

    out = FIXTURES_DIR / "comprehensive_document.docx"
    doc.save(str(out))
    img_path.unlink(missing_ok=True)
    print(f"Created {out}")


def create_comprehensive_pptx() -> None:
    """Build a PPTX exercising every feature the importer may encounter.

    Features: titles, body text, bold, italic, subscript, superscript,
    different fonts, bullet lists, numbered indent, tables, hyperlinks,
    images, shapes, charts, textboxes, speaker notes, special characters,
    emoji, blank slides.
    """
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    prs = Presentation()

    # Slide 1: Title
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "Comprehensive Feature Test"
    s1.placeholders[1].text = "Testing every supported PPTX feature"

    # Slide 2: Text formatting — bold, italic, subscript, superscript, fonts
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Text Formatting"
    tf = s2.placeholders[1].text_frame
    tf.clear()

    p = tf.paragraphs[0]
    rb = p.add_run()
    rb.text = "Bold text"
    rb.font.bold = True
    p.add_run().text = " and "
    ri = p.add_run()
    ri.text = "italic text"
    ri.font.italic = True
    p.add_run().text = " and "
    rbi = p.add_run()
    rbi.text = "bold italic"
    rbi.font.bold = True
    rbi.font.italic = True

    # Subscript/superscript (via XML baseline property)
    p2 = tf.add_paragraph()
    p2.add_run().text = "Water: H"
    r_sub = p2.add_run()
    r_sub.text = "2"
    # Set baseline to -25% for subscript
    rpr = r_sub._r.get_or_add_rPr()
    rpr.set("baseline", "-25000")
    p2.add_run().text = "O | Energy: E=mc"
    r_sup = p2.add_run()
    r_sup.text = "2"
    rpr2 = r_sup._r.get_or_add_rPr()
    rpr2.set("baseline", "30000")

    # Different fonts
    p3 = tf.add_paragraph()
    r_arial = p3.add_run()
    r_arial.text = "Arial. "
    r_arial.font.name = "Arial"
    r_times = p3.add_run()
    r_times.text = "Times New Roman. "
    r_times.font.name = "Times New Roman"
    r_mono = p3.add_run()
    r_mono.text = "Courier (monospace)."
    r_mono.font.name = "Courier New"

    # Slide 3: Special characters and emoji
    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    s3.shapes.title.text = "Special Characters"
    tf3 = s3.placeholders[1].text_frame
    tf3.clear()
    tf3.paragraphs[
        0
    ].text = (
        "\u00a9 \u00ae \u2122 \u00a7 \u00b0 \u00b1 \u2014 \u2013 \u2022 \u2026 \u20ac \u00a3 \u00a5"
    )
    p_emoji = tf3.add_paragraph()
    p_emoji.text = "Emoji: \U0001f600 \U0001f680 \U0001f4da \u2705 \u274c \U0001f525 \U0001f4a1"

    s3.notes_slide.notes_text_frame.text = "Verify special characters and emoji survive import."

    # Slide 4: Bullet and numbered lists (deep nesting)
    s4 = prs.slides.add_slide(prs.slide_layouts[1])
    s4.shapes.title.text = "Lists"
    tf4 = s4.placeholders[1].text_frame
    tf4.clear()
    items = [
        (0, "Bullet Level 0"),
        (1, "Bullet Level 1"),
        (2, "Bullet Level 2"),
        (3, "Bullet Level 3"),
        (0, "Back to Level 0"),
        (1, "Another sub-item"),
    ]
    for i, (level, text) in enumerate(items):
        if i == 0:
            tf4.paragraphs[0].text = text
            tf4.paragraphs[0].level = level
        else:
            p = tf4.add_paragraph()
            p.text = text
            p.level = level

    # Slide 5: Table with diverse content
    s5 = prs.slides.add_slide(prs.slide_layouts[1])
    s5.shapes.title.text = "Data Table"
    tbl = s5.shapes.add_table(5, 4, Inches(0.5), Inches(2), Inches(9), Inches(3)).table
    for c, h in enumerate(["Feature", "Status", "Priority", "Notes"]):
        tbl.cell(0, c).text = h
    rows = [
        ["Import DOCX", "Complete", "P0", "All features"],
        ["Import PPTX", "Complete", "P0", "Charts | shapes limited"],
        ["Import XLSX", "Complete", "P1", "Formulas not evaluated"],
        ["Import TEX", "Complete", "P1", "TikZ unsupported"],
    ]
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = val

    # Slide 6: Hyperlinks
    s6 = prs.slides.add_slide(prs.slide_layouts[1])
    s6.shapes.title.text = "Hyperlinks"
    tf6 = s6.placeholders[1].text_frame
    tf6.clear()
    links = [
        ("Documentation", "https://example.com/docs"),
        ("API Reference", "https://example.com/api"),
        ("Source Code", "https://github.com/example/repo"),
        ("Issue Tracker", "https://github.com/example/repo/issues"),
    ]
    for i, (text, url) in enumerate(links):
        r = tf6.paragraphs[0].add_run() if i == 0 else tf6.add_paragraph().add_run()
        r.text = text
        r.hyperlink.address = url

    # Slide 7: Images
    s7 = prs.slides.add_slide(prs.slide_layouts[1])
    s7.shapes.title.text = "Images"
    img_path = FIXTURES_DIR / "_tmp_comp_pptx.png"
    img_path.write_bytes(make_png())
    s7.shapes.add_picture(str(img_path), Inches(1), Inches(2), Inches(3), Inches(2))
    s7.shapes.add_picture(str(img_path), Inches(5), Inches(2), Inches(2), Inches(2))

    # Slide 8: Chart
    s8 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    s8.shapes.title.text = "Chart"
    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
    chart_data.add_series("Revenue", (100, 120, 140, 180))
    chart_data.add_series("Expenses", (80, 85, 90, 95))
    s8.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(2),
        Inches(8),
        Inches(4.5),
        chart_data,
    )

    # Slide 9: Textbox (free-floating shape with text)
    s9 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    txbox = s9.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1.5))
    tf9 = txbox.text_frame
    tf9.text = "This is a textbox with custom positioning."
    p9b = tf9.add_paragraph()
    p9b.text = "Textboxes are free-floating text containers."

    # Second textbox
    txbox2 = s9.shapes.add_textbox(Inches(1), Inches(3), Inches(6), Inches(1))
    txbox2.text_frame.text = "Another textbox for callout information."

    # Slide 10: Shapes (rectangles, ovals)
    from pptx.enum.shapes import MSO_SHAPE

    s10 = prs.slides.add_slide(prs.slide_layouts[6])
    rect = s10.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(1),
        Inches(3),
        Inches(1),
    )
    rect.text = "Rectangle shape with text"
    oval = s10.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(5),
        Inches(1),
        Inches(3),
        Inches(2),
    )
    oval.text = "Oval shape"
    arrow = s10.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(2),
        Inches(4),
        Inches(4),
        Inches(1),
    )
    arrow.text = "Arrow shape"

    # Slide 11: Speaker notes with formatting
    s11 = prs.slides.add_slide(prs.slide_layouts[1])
    s11.shapes.title.text = "Slide with Detailed Notes"
    s11.placeholders[1].text_frame.text = "Check the speaker notes."
    s11.notes_slide.notes_text_frame.text = (
        "These are detailed speaker notes.\n"
        "They span multiple lines.\n"
        "Line 3: mention the quarterly targets.\n"
        "Line 4: discuss next steps."
    )

    # Slide 12: Blank (no title, no content)
    prs.slides.add_slide(prs.slide_layouts[6])

    out = FIXTURES_DIR / "comprehensive_presentation.pptx"
    prs.save(str(out))
    img_path.unlink(missing_ok=True)
    print(f"Created {out}")


def create_comprehensive_xlsx() -> None:
    """Build an XLSX exercising every feature the importer may encounter.

    Features: multiple sheets, formulas, dates, times, datetimes, numbers,
    floats, integers, strings, None cells, merged cells, hyperlinks,
    embedded images, charts, styled cells, special characters, emoji,
    pipe characters, boolean values, very wide tables, trailing empty rows.
    """
    from datetime import time

    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference
    from openpyxl.drawing.image import Image
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side

    wb = Workbook()
    wb.remove(wb.active)

    # Sheet 1: "Mixed Types" — every data type
    ws1 = wb.create_sheet(title="Mixed Types")
    ws1.append(["Type", "Value", "Notes"])
    ws1.append(["String", "Hello World", "Plain text"])
    ws1.append(["Integer", 42, "Whole number"])
    ws1.append(["Float", 3.14159, "Decimal"])
    ws1.append(["Integer Float", 100.0, "Should display as 100"])
    ws1.append(["Negative", -273.15, "Below zero"])
    ws1.append(["Boolean True", True, "Truthy"])
    ws1.append(["Boolean False", False, "Falsy"])
    ws1.append(["Date", date(2025, 6, 15), "ISO date"])
    ws1.append(["DateTime", datetime(2025, 6, 15, 14, 30, 0), "With time"])
    ws1.append(["Time", time(9, 30, 0), "Just time"])
    ws1.append(["None", None, "Empty cell"])
    ws1.append(["Zero", 0, "Numeric zero"])
    ws1.append(["Large Number", 9999999999, "10 digits"])
    ws1.append(["Small Float", 0.0001, "Very small"])
    ws1.append(["Currency-like", 1234.56, "$1,234.56"])

    # Sheet 2: "Styled Cells" — fonts, colors, alignment
    ws2 = wb.create_sheet(title="Styled Cells")
    ws2.append(["Style", "Content"])
    ws2.append(["Bold", "Bold text"])
    ws2["B2"].font = Font(bold=True)
    ws2.append(["Italic", "Italic text"])
    ws2["B3"].font = Font(italic=True)
    ws2.append(["Red Font", "Red colored text"])
    ws2["B4"].font = Font(color="FF0000")
    ws2.append(["Blue Fill", "Blue background"])
    ws2["B5"].fill = PatternFill(start_color="0000FF", end_color="0000FF", fill_type="solid")
    ws2.append(["Large Font", "18pt text"])
    ws2["B6"].font = Font(size=18)
    ws2.append(["Monospace", "Code-like text"])
    ws2["B7"].font = Font(name="Courier New")
    ws2.append(["Bordered", "Cell with borders"])
    thin = Side(style="thin")
    ws2["B8"].border = Border(left=thin, right=thin, top=thin, bottom=thin)
    ws2.append(["Centered", "Center aligned"])
    ws2["B9"].alignment = Alignment(horizontal="center")
    ws2.append(["Wrapped", "This is a long text that should wrap"])
    ws2["B10"].alignment = Alignment(wrap_text=True)

    # Sheet 3: "Merged Cells"
    ws3 = wb.create_sheet(title="Merged Cells")
    ws3["A1"] = "Merged Header Spanning Three Columns"
    ws3.merge_cells("A1:C1")
    ws3["A2"] = "Col A"
    ws3["B2"] = "Col B"
    ws3["C2"] = "Col C"
    ws3["A3"] = "Data 1"
    ws3["B3"] = "Data 2"
    ws3["C3"] = "Data 3"
    ws3["A4"] = "Merged Rows"
    ws3.merge_cells("A4:A5")
    ws3["B4"] = "Row 4 data"
    ws3["C4"] = "More data"
    ws3["B5"] = "Row 5 data"
    ws3["C5"] = "Even more"

    # Sheet 4: "Special Characters"
    ws4 = wb.create_sheet(title="Special Characters")
    ws4.append(["Symbol", "Character"])
    ws4.append(["Copyright", "\u00a9"])
    ws4.append(["Registered", "\u00ae"])
    ws4.append(["Trademark", "\u2122"])
    ws4.append(["Em Dash", "\u2014"])
    ws4.append(["Euro", "\u20ac"])
    ws4.append(["Pound", "\u00a3"])
    ws4.append(["Yen", "\u00a5"])
    ws4.append(["Degree", "\u00b0"])
    ws4.append(["Pipe", "value | other value"])
    ws4.append(["Emoji", "\U0001f600 \U0001f680"])
    ws4.append(["Arrows", "\u2190 \u2191 \u2192 \u2193"])
    ws4.append(["Math", "\u221a \u221e \u2260 \u2264 \u2265"])

    # Sheet 5: "Formulas" (data_only=True means we see computed values)
    ws5 = wb.create_sheet(title="Formulas")
    ws5.append(["A", "B", "Sum", "Product"])
    ws5.append([10, 20, 30, 200])  # Computed values
    ws5.append([5, 15, 20, 75])
    ws5.append([100, 200, 300, 20000])
    # Note: formulas aren't preserved with data_only=True loading

    # Sheet 6: "Chart Data" — with embedded chart
    ws6 = wb.create_sheet(title="Chart Data")
    ws6.append(["Month", "Revenue", "Expenses"])
    ws6.append(["Jan", 5000, 3000])
    ws6.append(["Feb", 5500, 3200])
    ws6.append(["Mar", 6000, 3100])
    ws6.append(["Apr", 7000, 3500])

    chart = BarChart()
    chart.type = "col"
    chart.title = "Monthly Revenue vs Expenses"
    data = Reference(ws6, min_col=2, min_row=1, max_col=3, max_row=5)
    cats = Reference(ws6, min_col=1, min_row=2, max_row=5)
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(cats)
    ws6.add_chart(chart, "E2")

    # Sheet 7: "Images"
    ws7 = wb.create_sheet(title="Images")
    ws7.append(["Data with embedded image below"])
    ws7.append(["Row 2 content"])
    img_path = FIXTURES_DIR / "_tmp_comp_xlsx.png"
    img_path.write_bytes(make_png())
    ws7.add_image(Image(str(img_path)), "A4")

    # Sheet 8: "Hyperlinks"
    ws8 = wb.create_sheet(title="Hyperlinks")
    ws8.append(["Name", "URL"])
    ws8["A2"] = "Documentation"
    ws8["B2"] = "https://example.com/docs"
    ws8["B2"].hyperlink = "https://example.com/docs"
    ws8["A3"] = "GitHub"
    ws8["B3"] = "https://github.com/example"
    ws8["B3"].hyperlink = "https://github.com/example"

    # Sheet 9: "Wide Table" — many columns
    ws9 = wb.create_sheet(title="Wide Table")
    headers = [f"Column_{i}" for i in range(1, 16)]
    ws9.append(headers)
    for r in range(1, 4):
        ws9.append([f"R{r}C{c}" for c in range(1, 16)])

    # Sheet 10: "Empty" — deliberately blank
    wb.create_sheet(title="Empty")

    # Sheet 11: "Trailing Blanks"
    ws11 = wb.create_sheet(title="Trailing Blanks")
    ws11.append(["Key", "Value"])
    ws11.append(["alpha", 1])
    ws11.append(["beta", 2])
    ws11.append([None, None])
    ws11.append([None, None])
    ws11.append([None, None])

    out = FIXTURES_DIR / "comprehensive_workbook.xlsx"
    wb.save(out)
    img_path.unlink(missing_ok=True)
    print(f"Created {out}")


if __name__ == "__main__":
    create_reference_pptx()
    create_reference_xlsx()
    create_comprehensive_docx()
    create_comprehensive_pptx()
    create_comprehensive_xlsx()

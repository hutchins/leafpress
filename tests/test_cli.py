"""Tests for CLI commands."""

from pathlib import Path

from docx import Document as DocxDocument
from pptx import Presentation
from typer.testing import CliRunner

from leafpress.cli import cli

runner = CliRunner()


def test_version_flag() -> None:
    result = runner.invoke(cli, ["--version"])
    assert result.exit_code == 0
    assert "leafpress version" in result.output


def test_init_command(tmp_path: Path) -> None:
    result = runner.invoke(cli, ["init", str(tmp_path)])
    assert result.exit_code == 0
    assert (tmp_path / "leafpress.yml").exists()


def test_init_no_overwrite(tmp_path: Path) -> None:
    (tmp_path / "leafpress.yml").write_text("existing")
    result = runner.invoke(cli, ["init", str(tmp_path)])
    assert result.exit_code == 1
    assert "already exists" in result.output


def test_info_command(sample_mkdocs_dir: Path) -> None:
    result = runner.invoke(cli, ["info", str(sample_mkdocs_dir)])
    assert result.exit_code == 0
    assert "Test Documentation" in result.output


def test_convert_missing_source() -> None:
    result = runner.invoke(cli, ["convert", "/nonexistent/path"])
    assert result.exit_code == 1


def test_convert_pdf(sample_mkdocs_dir: Path, tmp_output: Path) -> None:
    result = runner.invoke(
        cli,
        ["convert", str(sample_mkdocs_dir), "-f", "pdf", "-o", str(tmp_output)],
    )
    assert result.exit_code == 0
    assert list(tmp_output.glob("*.pdf"))


def test_convert_with_watermark(sample_mkdocs_dir: Path, tmp_output: Path) -> None:
    result = runner.invoke(
        cli,
        [
            "convert",
            str(sample_mkdocs_dir),
            "-f",
            "pdf",
            "-o",
            str(tmp_output),
            "--watermark",
            "DRAFT",
        ],
    )
    assert result.exit_code == 0


def test_convert_with_local_time(sample_mkdocs_dir: Path, tmp_output: Path) -> None:
    result = runner.invoke(
        cli,
        [
            "convert",
            str(sample_mkdocs_dir),
            "-f",
            "pdf",
            "-o",
            str(tmp_output),
            "--local-time",
        ],
    )
    assert result.exit_code == 0


def test_convert_no_cover_no_toc(sample_mkdocs_dir: Path, tmp_output: Path) -> None:
    result = runner.invoke(
        cli,
        [
            "convert",
            str(sample_mkdocs_dir),
            "-f",
            "pdf",
            "-o",
            str(tmp_output),
            "--no-cover-page",
            "--no-toc",
        ],
    )
    assert result.exit_code == 0


def test_info_nonexistent_source() -> None:
    result = runner.invoke(cli, ["info", "/nonexistent/path"])
    assert result.exit_code == 1


# --- Multi-file import tests ---


def _make_docx(path: Path, title: str = "Test") -> Path:
    """Create a minimal DOCX file."""
    doc = DocxDocument()
    doc.add_heading(title, level=1)
    doc.save(str(path))
    return path


def _make_pptx(path: Path, title: str = "Slides") -> Path:
    """Create a minimal PPTX file."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    prs.save(str(path))
    return path


def test_import_multiple_docx(tmp_path: Path) -> None:
    """Import multiple DOCX files in one command."""
    f1 = _make_docx(tmp_path / "a.docx", "Doc A")
    f2 = _make_docx(tmp_path / "b.docx", "Doc B")
    out = tmp_path / "out"
    out.mkdir()

    result = runner.invoke(cli, ["import", str(f1), str(f2), "-o", str(out)])
    assert result.exit_code == 0
    assert (out / "a.md").exists()
    assert (out / "b.md").exists()
    assert result.output.count("Done!") == 2


def test_import_mixed_formats(tmp_path: Path) -> None:
    """Import a mix of DOCX and PPTX files."""
    docx = _make_docx(tmp_path / "report.docx", "Report")
    pptx = _make_pptx(tmp_path / "deck.pptx", "Deck")
    out = tmp_path / "out"
    out.mkdir()

    result = runner.invoke(cli, ["import", str(docx), str(pptx), "-o", str(out)])
    assert result.exit_code == 0
    assert (out / "report.md").exists()
    assert (out / "deck.md").exists()


def test_import_multiple_output_file_rejected(tmp_path: Path) -> None:
    """Using -o with a file path and multiple inputs is an error."""
    f1 = _make_docx(tmp_path / "a.docx")
    f2 = _make_docx(tmp_path / "b.docx")

    result = runner.invoke(cli, ["import", str(f1), str(f2), "-o", str(tmp_path / "single.md")])
    assert result.exit_code == 1
    assert "directory" in result.output.lower()


def test_import_partial_failure(tmp_path: Path) -> None:
    """One bad file doesn't stop other files from being imported."""
    good = _make_docx(tmp_path / "good.docx", "Good")
    bad = tmp_path / "missing.docx"
    out = tmp_path / "out"
    out.mkdir()

    result = runner.invoke(cli, ["import", str(good), str(bad), "-o", str(out)])
    assert result.exit_code == 1  # overall failure
    assert (out / "good.md").exists()  # first file succeeded
    assert "1 file(s) failed" in result.output


def test_import_single_file_still_works(sample_docx: Path, tmp_output: Path) -> None:
    """Single file argument still works (backwards compat)."""
    result = runner.invoke(cli, ["import", str(sample_docx), "-o", str(tmp_output)])
    assert result.exit_code == 0
    assert "Done!" in result.output


def test_convert_monorepo_config_no_mkdocs_required(tmp_path: Path) -> None:
    """convert -c with projects: should not require mkdocs.yml in cwd."""
    # Create a sub-project with mkdocs.yml
    sub = tmp_path / "services" / "api"
    sub.mkdir(parents=True)
    (sub / "mkdocs.yml").write_text("site_name: API Docs\nnav:\n  - Home: index.md\n")
    docs = sub / "docs"
    docs.mkdir()
    (docs / "index.md").write_text("# API\n\nHello.\n")

    # Create monorepo config in tmp_path (no mkdocs.yml here)
    config = tmp_path / "leafpress.yml"
    config.write_text(
        'company_name: "Test Corp"\nproject_name: "Monorepo Docs"\nprojects:\n  - services/api\n'
    )

    out = tmp_path / "output"
    out.mkdir()

    result = runner.invoke(
        cli,
        ["convert", "-c", str(config), "-f", "pdf", "-o", str(out)],
    )
    assert "Monorepo mode" in result.output
    assert result.exit_code == 0

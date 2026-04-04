"""Tests for PPTX to Markdown import converter."""

from __future__ import annotations

from pathlib import Path

import pytest
from helpers import make_png
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from leafpress.exceptions import PptxImportError
from leafpress.importer.converter_pptx import (
    _runs_to_markdown,
    _table_to_markdown,
    import_pptx,
)

# --- helpers ---


def _make_pptx(tmp_path: Path, slides: list[dict]) -> Path:
    """Create a PPTX file from slide definitions.

    Each slide dict can have:
        title: str
        body: str | list[str]
        notes: str
        table: list[list[str]]  (rows of cells)
    """
    prs = Presentation()
    for slide_data in slides:
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)

        if "title" in slide_data:
            slide.shapes.title.text = slide_data["title"]

        if "body" in slide_data:
            body_placeholder = slide.placeholders[1]
            tf = body_placeholder.text_frame
            tf.clear()
            body = slide_data["body"]
            if isinstance(body, str):
                body = [body]
            for i, text in enumerate(body):
                if i == 0:
                    tf.paragraphs[0].text = text
                else:
                    tf.add_paragraph().text = text

        if "notes" in slide_data:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["notes"]

        if "table" in slide_data:
            table_data = slide_data["table"]
            rows_count = len(table_data)
            cols_count = len(table_data[0]) if table_data else 0
            tbl = slide.shapes.add_table(
                rows_count,
                cols_count,
                Inches(1),
                Inches(3),
                Inches(6),
                Inches(2),
            ).table
            for r, row_data in enumerate(table_data):
                for c, cell_text in enumerate(row_data):
                    tbl.cell(r, c).text = cell_text

    path = tmp_path / "test.pptx"
    prs.save(path)
    return path


def _make_formatted_pptx(tmp_path: Path) -> Path:
    """Create a PPTX with bold/italic formatting."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Formatted"

    tf = slide.placeholders[1].text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run_bold = p.add_run()
    run_bold.text = "bold text"
    run_bold.font.bold = True
    run_normal = p.add_run()
    run_normal.text = " and "
    run_italic = p.add_run()
    run_italic.text = "italic text"
    run_italic.font.italic = True

    path = tmp_path / "formatted.pptx"
    prs.save(path)
    return path


def _make_image_pptx(tmp_path: Path) -> Path:
    """Create a PPTX with an embedded image."""
    img_path = tmp_path / "red.png"
    img_path.write_bytes(make_png())

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "With Image"
    slide.shapes.add_picture(str(img_path), Inches(1), Inches(2), Inches(2), Inches(2))

    path = tmp_path / "with_image.pptx"
    prs.save(path)
    return path


# --- import_pptx tests ---


def test_import_basic_slides(tmp_path: Path) -> None:
    """Basic slides with title and body text are converted."""
    pptx_path = _make_pptx(
        tmp_path,
        [
            {"title": "Introduction", "body": "Welcome to the presentation."},
            {"title": "Details", "body": "Here are the details."},
        ],
    )
    result = import_pptx(pptx_path, extract_images=False)
    md = result.markdown_path.read_text()
    assert "## Introduction" in md
    assert "Welcome to the presentation." in md
    assert "## Details" in md
    assert "Here are the details." in md


def test_slide_titles_as_h2(tmp_path: Path) -> None:
    """Each slide gets an H2 heading from its title."""
    pptx_path = _make_pptx(
        tmp_path,
        [
            {"title": "First Slide"},
            {"title": "Second Slide"},
        ],
    )
    result = import_pptx(pptx_path, extract_images=False)
    md = result.markdown_path.read_text()
    assert md.count("## ") == 2
    assert "## First Slide" in md
    assert "## Second Slide" in md


def test_untitled_slide_fallback(tmp_path: Path) -> None:
    """Slides without titles get '## Slide N' heading."""
    prs = Presentation()
    # Use blank layout (index 6) which has no title placeholder
    prs.slides.add_slide(prs.slide_layouts[6])
    path = tmp_path / "untitled.pptx"
    prs.save(path)

    result = import_pptx(path, extract_images=False)
    md = result.markdown_path.read_text()
    assert "## Slide 1" in md


def test_bold_italic_formatting(tmp_path: Path) -> None:
    """Bold and italic formatting is preserved."""
    pptx_path = _make_formatted_pptx(tmp_path)
    result = import_pptx(pptx_path, extract_images=False)
    md = result.markdown_path.read_text()
    assert "**bold text**" in md
    assert "*italic text*" in md


def test_multiple_body_paragraphs(tmp_path: Path) -> None:
    """Multiple body paragraphs are all extracted."""
    pptx_path = _make_pptx(
        tmp_path,
        [
            {"title": "Multi", "body": ["Line one", "Line two", "Line three"]},
        ],
    )
    result = import_pptx(pptx_path, extract_images=False)
    md = result.markdown_path.read_text()
    assert "Line one" in md
    assert "Line two" in md
    assert "Line three" in md


def test_table_extraction(tmp_path: Path) -> None:
    """Tables are converted to pipe-style markdown."""
    pptx_path = _make_pptx(
        tmp_path,
        [
            {
                "title": "Data",
                "table": [
                    ["Name", "Value"],
                    ["Alpha", "100"],
                    ["Beta", "200"],
                ],
            },
        ],
    )
    result = import_pptx(pptx_path, extract_images=False)
    md = result.markdown_path.read_text()
    assert "| Name" in md
    assert "| Alpha" in md
    assert "| Beta" in md
    assert "---" in md  # separator row


def test_image_extraction(tmp_path: Path) -> None:
    """Embedded images are extracted to assets/ and referenced in markdown."""
    pptx_path = _make_image_pptx(tmp_path)
    result = import_pptx(pptx_path)
    md = result.markdown_path.read_text()
    assert "![](assets/" in md
    assert len(result.images) == 1
    assert result.images[0].exists()


def test_speaker_notes_as_blockquote(tmp_path: Path) -> None:
    """Speaker notes are included as blockquotes."""
    pptx_path = _make_pptx(
        tmp_path,
        [
            {"title": "Noted", "notes": "Remember to explain this carefully."},
        ],
    )
    result = import_pptx(pptx_path, extract_images=False)
    md = result.markdown_path.read_text()
    assert "> Remember to explain this carefully." in md


def test_no_notes_flag(tmp_path: Path) -> None:
    """include_notes=False omits speaker notes."""
    pptx_path = _make_pptx(
        tmp_path,
        [
            {"title": "Noted", "notes": "Secret notes here."},
        ],
    )
    result = import_pptx(pptx_path, extract_images=False, include_notes=False)
    md = result.markdown_path.read_text()
    assert "Secret notes here" not in md


def test_empty_presentation(tmp_path: Path) -> None:
    """Empty presentation produces minimal output."""
    prs = Presentation()
    path = tmp_path / "empty.pptx"
    prs.save(path)
    result = import_pptx(path, extract_images=False)
    md = result.markdown_path.read_text()
    assert md.strip() == ""


def test_output_path_default(tmp_path: Path) -> None:
    """Default output uses same stem as input."""
    pptx_path = _make_pptx(tmp_path, [{"title": "Test"}])
    result = import_pptx(pptx_path, extract_images=False)
    assert result.markdown_path == tmp_path / "test.md"


def test_output_path_explicit(tmp_path: Path) -> None:
    """Explicit output path is respected."""
    pptx_path = _make_pptx(tmp_path, [{"title": "Test"}])
    out = tmp_path / "custom" / "output.md"
    result = import_pptx(pptx_path, output_path=out, extract_images=False)
    assert result.markdown_path == out
    assert out.exists()


def test_output_path_directory(tmp_path: Path) -> None:
    """Directory output creates <stem>.md inside it."""
    pptx_path = _make_pptx(tmp_path, [{"title": "Test"}])
    out_dir = tmp_path / "out"
    out_dir.mkdir()
    result = import_pptx(pptx_path, output_path=out_dir, extract_images=False)
    assert result.markdown_path == out_dir / "test.md"


def test_invalid_file_raises(tmp_path: Path) -> None:
    """Wrong extension raises PptxImportError."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("hello")
    with pytest.raises(PptxImportError, match=r"Not a \.pptx file"):
        import_pptx(txt_file)


def test_missing_file_raises(tmp_path: Path) -> None:
    """Nonexistent file raises PptxImportError."""
    with pytest.raises(PptxImportError, match="File not found"):
        import_pptx(tmp_path / "nonexistent.pptx")


# --- unit tests for helpers ---


def test_table_to_markdown_pipe_format() -> None:
    """_table_to_markdown produces correct pipe table."""
    from unittest.mock import MagicMock

    table = MagicMock()
    row1 = MagicMock()
    row1.cells = [MagicMock(text="H1"), MagicMock(text="H2")]
    row2 = MagicMock()
    row2.cells = [MagicMock(text="a"), MagicMock(text="b")]
    table.rows = [row1, row2]

    md = _table_to_markdown(table)
    assert "| H1" in md
    assert "| a" in md
    assert "---" in md


def test_runs_to_markdown_plain() -> None:
    """Plain text runs are joined."""
    from unittest.mock import MagicMock

    para = MagicMock()
    run = MagicMock()
    run.text = "hello world"
    run.font.bold = False
    run.font.italic = False
    run.hyperlink = None
    para.runs = [run]

    assert _runs_to_markdown(para) == "hello world"


def test_convert_shape_warns_on_chart() -> None:
    """Chart shapes produce a warning."""
    from unittest.mock import MagicMock

    from leafpress.importer.converter_pptx import _convert_shape

    shape = MagicMock()
    shape.shape_type = MSO_SHAPE_TYPE.CHART
    shape.name = "Chart 1"
    shape.has_text_frame = False
    warnings: list[str] = []

    result = _convert_shape(shape, None, "Revenue Slide", warnings)
    assert result == ""
    assert len(warnings) == 1
    assert "chart" in warnings[0].lower()
    assert "Chart 1" in warnings[0]
    assert "Revenue Slide" in warnings[0]
    assert "skipped" in warnings[0].lower()


def test_convert_shape_no_warn_on_freeform() -> None:
    """Freeform shapes (decorative lines, etc.) are silently skipped."""
    from unittest.mock import MagicMock

    from leafpress.importer.converter_pptx import _convert_shape

    shape = MagicMock()
    shape.shape_type = MSO_SHAPE_TYPE.FREEFORM
    shape.has_text_frame = False
    warnings: list[str] = []

    _convert_shape(shape, None, "Slide 1", warnings)
    assert len(warnings) == 0


def test_runs_to_markdown_hyperlink() -> None:
    """Hyperlinks in runs produce markdown link syntax."""
    from unittest.mock import MagicMock

    para = MagicMock()
    run = MagicMock()
    run.text = "Click here"
    run.font.bold = False
    run.font.italic = False
    run.hyperlink.address = "https://example.com"
    para.runs = [run]

    md = _runs_to_markdown(para)
    assert "[Click here](https://example.com)" in md


def test_runs_to_markdown_bold_hyperlink() -> None:
    """Bold text with a hyperlink unwraps formatting for the link text."""
    from unittest.mock import MagicMock

    para = MagicMock()
    run = MagicMock()
    run.text = "bold link"
    run.font.bold = True
    run.font.italic = False
    run.hyperlink.address = "https://example.com"
    para.runs = [run]

    md = _runs_to_markdown(para)
    assert "[bold link](https://example.com)" in md


def test_runs_to_markdown_empty_run() -> None:
    """Runs with empty text are skipped."""
    from unittest.mock import MagicMock

    para = MagicMock()
    empty_run = MagicMock()
    empty_run.text = ""
    text_run = MagicMock()
    text_run.text = "visible"
    text_run.font.bold = False
    text_run.font.italic = False
    text_run.hyperlink = None
    para.runs = [empty_run, text_run]

    md = _runs_to_markdown(para)
    assert md == "visible"


# --- CLI integration ---


def test_cli_pptx_integration(tmp_path: Path) -> None:
    """End-to-end CLI import of a PPTX file."""
    from typer.testing import CliRunner

    from leafpress.cli import cli

    pptx_path = _make_pptx(
        tmp_path,
        [
            {"title": "CLI Test", "body": "Content here."},
        ],
    )
    runner = CliRunner()
    result = runner.invoke(cli, ["import", str(pptx_path), "-o", str(tmp_path / "out.md")])
    assert result.exit_code == 0
    assert "Done!" in result.output
    assert (tmp_path / "out.md").exists()
    md = (tmp_path / "out.md").read_text()
    assert "## CLI Test" in md


def test_cli_unsupported_format(tmp_path: Path) -> None:
    """CLI rejects unsupported file types."""
    from typer.testing import CliRunner

    from leafpress.cli import cli

    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("hello")
    runner = CliRunner()
    result = runner.invoke(cli, ["import", str(txt_file)])
    assert result.exit_code == 1
    assert "Unsupported file type" in result.output


# ---------------------------------------------------------------------------
# Comprehensive PPTX fixture — realistic multi-slide presentation
# ---------------------------------------------------------------------------


def _make_comprehensive_pptx(tmp_path: Path) -> Path:
    """Create a multi-slide PPTX with mixed content types."""
    prs = Presentation()

    # Slide 1: Title + formatted body text
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "Project Status"
    tf = slide1.placeholders[1].text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run_b = p.add_run()
    run_b.text = "On track"
    run_b.font.bold = True
    run_n = p.add_run()
    run_n.text = " for Q4 delivery with "
    run_i = p.add_run()
    run_i.text = "minor risks"
    run_i.font.italic = True

    # Slide 2: Table + speaker notes
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Budget Overview"
    tbl = slide2.shapes.add_table(3, 3, Inches(1), Inches(2), Inches(6), Inches(2)).table
    for c, h in enumerate(["Category", "Budgeted", "Actual"]):
        tbl.cell(0, c).text = h
    for c, v in enumerate(["Engineering", "$500K", "$480K"]):
        tbl.cell(1, c).text = v
    for c, v in enumerate(["Marketing", "$200K", "$220K"]):
        tbl.cell(2, c).text = v
    notes2 = slide2.notes_slide
    notes2.notes_text_frame.text = "Emphasize that engineering is under budget."

    # Slide 3: Indented text (multiple levels)
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Milestones"
    tf3 = slide3.placeholders[1].text_frame
    tf3.clear()
    tf3.paragraphs[0].text = "Phase 1: Foundation"
    p1 = tf3.add_paragraph()
    p1.text = "Database schema finalized"
    p1.level = 1
    p2 = tf3.add_paragraph()
    p2.text = "API contracts signed off"
    p2.level = 1
    p3 = tf3.add_paragraph()
    p3.text = "Phase 2: Build"
    p3.level = 0
    p4 = tf3.add_paragraph()
    p4.text = "Core services deployed"
    p4.level = 1
    p5 = tf3.add_paragraph()
    p5.text = "Integration testing complete"
    p5.level = 2

    # Slide 4: Image
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Architecture Diagram"
    img_path = tmp_path / "diagram.png"
    img_path.write_bytes(make_png())
    slide4.shapes.add_picture(str(img_path), Inches(1), Inches(2), Inches(4), Inches(3))

    # Slide 5: Blank layout (no title placeholder)
    prs.slides.add_slide(prs.slide_layouts[6])

    path = tmp_path / "comprehensive.pptx"
    prs.save(path)
    return path


class TestComprehensivePptx:
    """Tests using a realistic multi-slide presentation."""

    @pytest.fixture(autouse=True)
    def _convert(self, tmp_path: Path) -> None:
        pptx_path = _make_comprehensive_pptx(tmp_path)
        out_dir = tmp_path / "output"
        out_dir.mkdir()
        self.result = import_pptx(pptx_path, output_path=out_dir)
        self.content = self.result.markdown_path.read_text()

    def test_converts_successfully(self) -> None:
        """All slides convert without errors."""
        assert self.result.markdown_path.exists()
        assert len(self.content) > 200

    def test_slide_titles_as_headings(self) -> None:
        """Titled slides get ## headings."""
        assert "## Project Status" in self.content
        assert "## Budget Overview" in self.content
        assert "## Milestones" in self.content
        assert "## Architecture Diagram" in self.content

    def test_untitled_slide_fallback(self) -> None:
        """Blank-layout slide gets ## Slide N heading."""
        assert "## Slide 5" in self.content

    def test_bold_formatting(self) -> None:
        """Bold text in slide body preserved."""
        assert "**On track**" in self.content

    def test_italic_formatting(self) -> None:
        """Italic text in slide body preserved."""
        assert "*minor risks*" in self.content

    def test_table(self) -> None:
        """Table converted to pipe-style markdown."""
        assert "| Category" in self.content
        assert "Engineering" in self.content
        assert "$480K" in self.content
        assert "---" in self.content

    def test_speaker_notes(self) -> None:
        """Speaker notes rendered as blockquotes."""
        assert "> Emphasize that engineering is under budget." in self.content

    def test_indented_text(self) -> None:
        """Indented paragraphs produce nested bullets."""
        assert "Phase 1: Foundation" in self.content
        assert "  - Database schema finalized" in self.content
        assert "  - API contracts signed off" in self.content

    def test_deeper_indent(self) -> None:
        """Level-2 indent produces double-indented bullet."""
        assert "    - Integration testing complete" in self.content

    def test_image_extracted(self) -> None:
        """Embedded image extracted to assets/."""
        assert len(self.result.images) >= 1
        assert self.result.images[0].exists()
        assert "![](assets/" in self.content

    def test_slide_count(self) -> None:
        """All 5 slides appear in output."""
        assert self.content.count("## ") == 5
